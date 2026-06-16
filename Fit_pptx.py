import tkinter as tk
from tkinterdnd2 import DND_FILES, TkinterDnD
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

def drop(event):
    file_path = event.data.strip("{}")
    print("DROP:", repr(file_path))
    create_pptx(file_path)

def create_pptx(img_path):
    # スクリプトのあるフォルダ
    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "output.pptx")

    # 画像読み込み
    img = Image.open(img_path)
    w, h = img.size  # px

    # px → inch（96dpi換算）
    width_inch = w / 96
    height_inch = h / 96

    # pptx 作成
    prs = Presentation()

    # ★ スライドサイズを画像と同じにする
    prs.slide_width = Inches(width_inch)
    prs.slide_height = Inches(height_inch)

    # 白紙スライド
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 画像をスライド左上に等倍で貼る
    slide.shapes.add_picture(
        img_path,
        Inches(0), Inches(0),
        width=Inches(width_inch),
        height=Inches(height_inch)
    )

    prs.save(out_path)
    print("作成完了:", out_path)

# GUI
root = TkinterDnD.Tk()
root.title("画像 → PPTX 変換ツール")
root.configure(bg="#e34841")  # ← 好きな色に変更
root.geometry("400x200")

label = tk.Label(root, text="ここに画像をドロップ", width=40, height=10)
label.pack()

label.drop_target_register(DND_FILES)
label.dnd_bind('<<Drop>>', drop)

root.mainloop()
